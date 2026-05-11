"""
Real-Time Facial Emotion Recognition — Master-level defense presentation
13 slides · Dark navy/cyan AI theme · 15-minute presentation
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE  = Path('/home/charbel-mezeraani/cv project')
PLOTS = BASE / 'evaluation' / 'plots'
OUT   = BASE / 'Emotion_Recognition_Presentation.pptx'

BG     = RGBColor(0x0D, 0x11, 0x17)
CARD   = RGBColor(0x16, 0x1B, 0x22)
CARD2  = RGBColor(0x1C, 0x21, 0x28)
BLUE   = RGBColor(0x58, 0xA6, 0xFF)
CYAN   = RGBColor(0x39, 0xC5, 0xCF)
WHITE  = RGBColor(0xE6, 0xED, 0xF3)
GRAY   = RGBColor(0x8B, 0x94, 0x9E)
LGRAY  = RGBColor(0x48, 0x4F, 0x58)
GOLD   = RGBColor(0xD2, 0x99, 0x22)
GREEN  = RGBColor(0x3F, 0xB9, 0x50)
RED    = RGBColor(0xF8, 0x51, 0x49)
ORANGE = RGBColor(0xE3, 0x6B, 0x1B)
BORDER = RGBColor(0x30, 0x36, 0x3D)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def S():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill  = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def R(s, l, t, w, h, color=CARD, border=None, bw=0.75, rnd=False):
    sh = s.shapes.add_shape(5 if rnd else 1,
                            Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(bw)
    else:
        sh.line.fill.background()
    return sh

def T(s, text, l, t, w, h, sz=16, bold=False, color=WHITE,
      align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(sz); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color; run.font.name = 'Calibri'
    return txb

def T2(s, lines, l, t, w, h, align=PP_ALIGN.LEFT):
    txb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame; tf.word_wrap = True
    for i, (text, sz, bold, color, italic) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; run = p.add_run()
        run.text = text; run.font.size = Pt(sz); run.font.bold = bold
        run.font.italic = italic; run.font.color.rgb = color; run.font.name = 'Calibri'
    return txb

def IMG(s, path, l, t, w, h=None):
    p = Path(path)
    if p.exists():
        try:
            if h: return s.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))
            return s.shapes.add_picture(str(p), Inches(l), Inches(t), width=Inches(w))
        except: pass
    R(s, l, t, w, h or 2.0, color=LGRAY, border=BORDER)
    T(s, f'[{p.name}]', l+0.1, t+(h or 2.0)/2-0.2, w-0.2, 0.4, sz=9, color=GRAY, align=PP_ALIGN.CENTER)

def HDR(s, title, sub=None):
    R(s, 0, 0, 13.33, 0.88, color=CARD)
    R(s, 0, 0, 0.07,  0.88, color=BLUE)
    T(s, title, 0.22, 0.04, 10.5, 0.55, sz=26, bold=True, color=WHITE)
    if sub: T(s, sub, 0.22, 0.56, 10.0, 0.30, sz=12, color=GRAY)

def CHIP(s, label, l, t, w=1.7, h=0.42, bg=CARD2, fg=BLUE, sz=11):
    R(s, l, t, w, h, color=bg, border=fg, bw=0.75, rnd=True)
    T(s, label, l+0.05, t+0.03, w-0.1, h-0.06, sz=sz, bold=True, color=fg, align=PP_ALIGN.CENTER)

def NOTES(s, text):
    s.notes_slide.notes_text_frame.text = text

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════
s1 = S()
R(s1, 7.8, 0, 5.53, 7.5, color=CARD)
R(s1, 7.8, 0, 0.05, 7.5, color=BLUE)
T(s1, 'AI  ·  Computer Vision  ·  Master 1', 0.4, 0.5, 7.0, 0.45, sz=12, italic=True, color=GRAY)
T2(s1, [('Real-Time Facial', 34, True, WHITE, False),
        ('Emotion Recognition', 34, True, BLUE, False)], 0.4, 1.05, 7.15, 1.6)
T(s1, 'From Baseline to EfficientNetB2', 0.4, 2.72, 7.0, 0.5, sz=17, color=GRAY)
R(s1, 0.4, 3.38, 6.8, 0.04, color=BLUE)
for row, (lbl, val) in enumerate([('Course','Artificial Intelligence & Computer Vision'),
                                   ('Student','[Student Name(s)]'),
                                   ('University','[University Name]'),
                                   ('Date','May 2026')]):
    T(s1, lbl+':',  0.4,  3.55+row*0.42, 1.3, 0.38, sz=11, bold=True,  color=BLUE)
    T(s1, val,      1.75, 3.55+row*0.42, 5.8, 0.38, sz=11, color=WHITE)
T(s1, 'RAF-DB · 3,068 test images · RTX 4050 Laptop GPU · TensorFlow 2.21',
  0.4, 7.05, 7.1, 0.35, sz=9, italic=True, color=LGRAY)
T(s1, 'Best Result',  8.2, 0.9,  4.8, 0.5,  sz=14, color=GRAY, align=PP_ALIGN.CENTER)
T(s1, '80.02%',       8.0, 1.35, 5.1, 1.2,  sz=62, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
T(s1, 'Test Accuracy (RAF-DB)', 8.0, 2.55, 5.1, 0.4, sz=12, color=GRAY, align=PP_ALIGN.CENTER)
R(s1, 8.3, 3.1, 4.5, 0.04, color=BORDER)
T(s1, '+27.22%',             8.0, 3.25, 5.1, 0.75, sz=34, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
T(s1, 'vs DeepFace zero-shot', 8.0, 4.0,  5.1, 0.4,  sz=11, color=GRAY, align=PP_ALIGN.CENTER)
R(s1, 8.3, 4.55, 4.5, 0.04, color=BORDER)
ems = ['Angry','Disgust','Fear','Happy','Sad','Surprise','Neutral']
for i, em in enumerate(ems):
    c = i%2; r = i//2
    ex = 8.2+c*2.45; ey = 4.75+r*0.52
    if i==6: ex=9.4; ey=4.75+3*0.52
    CHIP(s1, em, ex, ey, w=2.2, h=0.4, sz=10)
NOTES(s1, "Welcome. This project builds a Real-Time Facial Emotion Recognition system, "
     "progressing from a 52.80% zero-shot DeepFace baseline to 80.02% with EfficientNetB2 — "
     "a +27.22% improvement. I'll cover: dataset, preprocessing, 6 model configurations, "
     "per-class analysis, deployment, and lessons learned. About 15 minutes total.")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM AND OBJECTIVES
# ═══════════════════════════════════════════════════════════════════
s2 = S()
HDR(s2, 'Why Is Emotion Recognition Hard?', 'Three core challenges — and four objectives')
challenges = [
    ('Class Imbalance', 'Happy & Neutral dominate naturalistic data. Disgust & Fear are rare — the model biases toward majority classes without correction.', RED),
    ('Visual Ambiguity', 'Fear and Surprise share identical primary muscles. Low-intensity anger resembles Neutral. Subtle cues require fine-grained discrimination.', ORANGE),
    ('Real-World Variation', 'Illumination, head pose, partial occlusion, and cultural norms all degrade recognition. Benchmark accuracy overstates real-world performance.', BLUE),
]
for i, (title, body, col) in enumerate(challenges):
    cx = 0.35+i*4.3
    R(s2, cx, 1.0, 4.05, 2.65, color=CARD, border=col, bw=1.5, rnd=True)
    R(s2, cx, 1.0, 4.05, 0.44, color=CARD2)
    T(s2, title, cx+0.14, 1.05, 3.75, 0.34, sz=14, bold=True, color=col)
    T(s2, body,  cx+0.14, 1.52, 3.75, 2.0,  sz=11.5, color=WHITE, wrap=True)
T(s2, '7 Target Emotion Classes:', 0.35, 3.88, 4.0, 0.38, sz=13, bold=True, color=GRAY)
edata = [('Angry',RED),('Disgust',ORANGE),('Fear',ORANGE),('Happy',GREEN),('Sad',BLUE),('Surprise',BLUE),('Neutral',GRAY)]
for i,(em,col) in enumerate(edata):
    CHIP(s2, em, 0.35+i*1.87, 4.32, w=1.75, h=0.48, fg=col, sz=11)
R(s2, 0.35, 5.05, 12.6, 1.15, color=CARD2, border=BLUE, bw=1.0, rnd=True)
T(s2, '"An 80% accuracy system can still misclassify the majority of disgust and fear expressions. Overall accuracy is insufficient."',
  0.6, 5.12, 12.1, 0.6, sz=13.5, italic=True, color=WHITE)
T(s2, 'Per-class analysis is mandatory for any real FER deployment.',
  0.6, 5.67, 12.1, 0.4, sz=12, bold=True, color=BLUE)
T(s2, 'Objectives:  1. Measure baseline     2. Train & compare 6 models     3. Deploy real-time system     4. Per-class analysis',
  0.35, 6.35, 12.6, 0.42, sz=11.5, color=GRAY)
NOTES(s2, "FER is hard for three reasons: class imbalance (disgust rarely appears in naturalistic "
     "photography), visual ambiguity (fear and surprise look almost identical in still images), "
     "and real-world variation. My 4 objectives: establish a zero-shot baseline, train and "
     "compare 6 models, analyze per-class performance, and deploy a real-time system. "
     "The key insight: 80% overall accuracy can coexist with F1=0.42 for disgust.")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — DATASET AND PREPROCESSING
# ═══════════════════════════════════════════════════════════════════
s3 = S()
HDR(s3, 'RAF-DB Dataset & Preprocessing', 'Real-world unconstrained images · critical engineering discovery')
R(s3, 0.35, 1.0, 4.1, 5.9, color=CARD, border=BORDER, rnd=True)
T(s3, 'RAF-DB Dataset', 0.55, 1.12, 3.7, 0.4, sz=14, bold=True, color=BLUE)
dstats = [('Test set','3,068 images'),('Training','RAF-DB + FER2013 + Extra'),
          ('After oversampling','84,000 images (12K/class)'),('Resolution','100×100 px aligned'),
          ('Classes','7 basic emotions'),('Hardest class','Disgust (very rare)')]
for i,(label,val) in enumerate(dstats):
    T(s3, label, 0.55, 1.65+i*0.72, 1.9, 0.42, sz=11, color=GRAY)
    T(s3, val,   2.5,  1.65+i*0.72, 1.75,0.42, sz=11, bold=True, color=WHITE)
    if i<5: R(s3, 0.55, 2.08+i*0.72, 3.7, 0.02, color=BORDER)
R(s3, 4.65, 1.0, 4.35, 5.9, color=CARD, border=BORDER, rnd=True)
T(s3, 'Preprocessing', 4.85, 1.12, 3.95, 0.4, sz=14, bold=True, color=BLUE)
T(s3, 'SE-CNN',         5.4, 1.65, 1.6, 0.35, sz=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
T(s3, 'EfficientNetB2', 7.0, 1.65, 1.7, 0.35, sz=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
prows = [('Input','48×48 px','96×96 px'),('Color','Grayscale','RGB'),
         ('Standardize','Per-image','Per-image'),('Augmentation','Geometric','Geometric'),
         ('Photometric','X  Excluded','X  Excluded'),('Oversampling','12K/class','Batch-level')]
for i,(p,v1,v2) in enumerate(prows):
    y=2.1+i*0.72; bg=CARD2 if i%2==0 else CARD
    R(s3, 4.65, y, 4.35, 0.68, color=bg)
    T(s3, p,  4.78, y+0.15, 1.55, 0.4, sz=10,   color=GRAY)
    T(s3, v1, 5.55, y+0.15, 1.55, 0.4, sz=10.5, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    T(s3, v2, 7.05, y+0.15, 1.65, 0.4, sz=10.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
R(s3, 9.2, 1.0, 3.8, 5.9, color=CARD, border=RED, bw=1.5, rnd=True)
T(s3, 'Critical Discovery', 9.38, 1.1, 3.45, 0.4, sz=13, bold=True, color=RED)
R(s3, 9.2, 1.5, 3.8, 0.03, color=RED)
T(s3, 'Photometric Augment + Per-Image Standardisation = Training Collapse',
  9.38, 1.6, 3.45, 0.9, sz=11.5, bold=True, color=WHITE, wrap=True)
T(s3, 'Accuracy frozen at 14.3% (random baseline) for entire 100-epoch run.',
  9.38, 2.55, 3.45, 0.75, sz=12, bold=True, color=RED, wrap=True)
T(s3, 'Root cause:', 9.38, 3.38, 3.45, 0.32, sz=11, bold=True, color=GRAY)
T(s3, 'Standardisation removes all brightness/contrast. RandomBrightness on top creates contradictory gradients. Model cannot converge.',
  9.38, 3.72, 3.45, 1.1, sz=10.5, color=WHITE, wrap=True)
T(s3, 'Fix: Geometric augmentation only\n(flip, rotate, zoom, translate)',
  9.38, 4.95, 3.45, 0.7, sz=12, bold=True, color=GREEN, wrap=True)
T(s3, 'Lesson: preprocessing and augmentation interact — test before full training.',
  9.38, 5.78, 3.45, 0.9, sz=10, italic=True, color=GRAY, wrap=True)
NOTES(s3, "RAF-DB is a real-world unconstrained dataset. After combining sources and oversampling, "
     "84,000 training images. A critical discovery: combining RandomBrightness/Contrast with "
     "per-image standardization caused training to freeze at 14.3% for the entire run. "
     "Fix: geometric augmentation only. This is not documented in most tutorials.")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM PIPELINE
# ═══════════════════════════════════════════════════════════════════
s4 = S()
HDR(s4, 'Complete System Pipeline', 'End-to-end: raw input → face detection → inference → Grad-CAM → deployment')
IMG(s4, PLOTS/'figure1_pipeline.png', 0.6, 0.95, 12.1)
R(s4, 0, 6.85, 13.33, 0.65, color=CARD)
for i,f in enumerate(['Haar Cascade  or  MTCNN face detection',
                       '8-pass TTA offline  ·  single-pass real-time',
                       'Grad-CAM integrated at every output mode']):
    T(s4, 'x  '+f, 0.4+i*4.45, 6.9, 4.2, 0.42, sz=10.5, color=GRAY)
NOTES(s4, "Full pipeline: image/video in → face detection (Haar or MTCNN) → crop & align → "
     "preprocess (resize, standardize) → model inference → 7-class softmax → Grad-CAM heatmap → "
     "5 output modes. For benchmarks: 8-pass TTA. For real-time: single-pass. "
     "Grad-CAM runs in all modes simultaneously.")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — MODEL PROGRESSION
# ═══════════════════════════════════════════════════════════════════
s5 = S()
HDR(s5, '6 Models — 1 Clear Winner', '52.80% → 80.02% through a principled progression')
IMG(s5, PLOTS/'pptx_model_progression.png', 0.35, 0.95, 12.6)
R(s5, 9.5, 5.9, 3.5, 1.38, color=CARD2, border=GOLD, bw=1.2, rnd=True)
T(s5, 'EfficientNetB2 vs Ensemble', 9.65, 5.95, 3.2, 0.38, sz=11, bold=True, color=GOLD)
T(s5, '+9.57% from one architectural change\n14M ImageNet images vs 84K custom training',
  9.65, 6.35, 3.2, 0.75, sz=11, color=WHITE, wrap=True)
NOTES(s5, "The arc arrow shows the +27.22% story from DeepFace to EfficientNetB2. Key transitions: "
     "+15.39% just from task-specific fine-tuning (MobileNetV2). Custom SE-CNN adds +1.54%. "
     "4-model ensemble adds +0.72%. Then EfficientNetB2 jumps +9.57% over the ensemble. "
     "The critical insight: pretrained features from 14M images cannot be replicated by "
     "architecture design at academic dataset scales.")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — SE-CNN ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════
s6 = S()
HDR(s6, 'Custom Multi-Scale SE-CNN', 'Purpose-built for FER spatial diversity — 3.2M params trained from scratch')
IMG(s6, PLOTS/'figure2_secnn.png', 0.3, 0.95, 8.3)
R(s6, 8.85, 0.95, 4.15, 6.35, color=CARD, border=BORDER, rnd=True)
T(s6, 'Architecture Details', 9.0, 1.05, 3.8, 0.38, sz=13, bold=True, color=BLUE)
R(s6, 8.85, 1.43, 4.15, 0.03, color=BORDER)
astats = [('Parameters','~3.2M',WHITE),('Input','48x48 grayscale',WHITE),
          ('Branches','3x3  5x5  7x7',CYAN),('SE Attention','Channel recalib.',CYAN),
          ('Residuals','Prevent vanishing grad',WHITE),('He Init + BN','Every conv layer',WHITE)]
for i,(k,v,vc) in enumerate(astats):
    y=1.52+i*0.62
    R(s6, 9.0, y, 3.85, 0.55, color=CARD2, rnd=True)
    T(s6, k, 9.12, y+0.08, 1.6,  0.38, sz=10.5, bold=True, color=GRAY)
    T(s6, v, 10.8, y+0.08, 1.9,  0.38, sz=10.5, bold=True, color=vc, align=PP_ALIGN.RIGHT)
R(s6, 8.85, 5.22, 4.15, 0.03, color=BORDER)
T(s6, 'SE-CNN v2 + SWA + TTA', 9.0, 5.3, 3.8, 0.35, sz=11.5, bold=True, color=WHITE)
T(s6, '69.73%', 9.0, 5.65, 3.8, 0.78, sz=40, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
T(s6, '+16.93% vs DeepFace', 9.0, 6.45, 3.8, 0.35, sz=11, color=GRAY, align=PP_ALIGN.CENTER)
NOTES(s6, "The SE-CNN insight: expression cues exist at multiple scales. Nasal wrinkle (disgust) "
     "needs a 3x3 kernel; a full smile needs 7x7. Three parallel branches capture all scales, "
     "then concatenated. SE attention adaptively suppresses identity/lighting channels and "
     "emphasizes expression-relevant ones. Despite being trained from scratch, reaches 69.73%.")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — EFFICIENTNETB2 FINE-TUNING
# ═══════════════════════════════════════════════════════════════════
s7 = S()
HDR(s7, 'EfficientNetB2 — Progressive Fine-Tuning', '4-phase strategy to prevent catastrophic forgetting · ImageNet compound-scaled backbone')
IMG(s7, PLOTS/'pptx_efficientnet_phases.png', 0.3, 0.92, 12.7)
for xp,txt in [('0.3','Warm up head\nbackbone frozen'),('3.3','Top-40% unlock\nbiggest jump +15.3%'),
                ('6.3','Full backbone\ngentle adaptation'),('9.2','SWA flat-minima\nbetter generalise'),
                ('11.7','8-pass TTA\n+4.86% final')]:
    T(s7, txt, float(xp), 6.82, 2.6, 0.55, sz=9.5, color=GRAY, align=PP_ALIGN.CENTER, wrap=True)
NOTES(s7, "Progressive unfreezing prevents catastrophic forgetting. Phase 1: freeze backbone, "
     "train head only (15 epochs). Phase 2: unfreeze top 40% — deepest, most task-specific — "
     "biggest jump 57->72%. Phase 3: full backbone at very low LR. Phase 4: SWA averages "
     "flat-minima weights. Final +4.86% from 8-pass TTA. Total: 80.02%.")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS
# ═══════════════════════════════════════════════════════════════════
s8 = S()
HDR(s8, 'Results: 80.02% Test Accuracy', 'RAF-DB · 3,068 images · 8-pass TTA · all results reproducible')
R(s8, 9.8, 0.95, 3.15, 1.55, color=CARD2, border=GOLD, bw=1.5, rnd=True)
T(s8, '80.02%', 9.82, 1.02, 3.1, 0.85, sz=46, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
T(s8, 'EfficientNetB2 + TTA', 9.82, 1.87, 3.1, 0.42, sz=11, color=GRAY, align=PP_ALIGN.CENTER)
R(s8, 9.8, 2.65, 3.15, 1.0, color=CARD2, border=GREEN, bw=1.2, rnd=True)
T(s8, '+27.22%', 9.82, 2.72, 3.1, 0.58, sz=30, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
T(s8, 'vs DeepFace baseline', 9.82, 3.32, 3.1, 0.28, sz=10, color=GRAY, align=PP_ALIGN.CENTER)
R(s8, 9.8, 3.78, 3.15, 0.85, color=CARD2, border=BLUE, bw=1.0, rnd=True)
T(s8, '+9.57%', 9.82, 3.83, 3.1, 0.5, sz=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
T(s8, 'vs 4-model ensemble', 9.82, 4.35, 3.1, 0.25, sz=10, color=GRAY, align=PP_ALIGN.CENTER)
IMG(s8, PLOTS/'model_comparison.png', 0.3, 0.95, 9.3)
R(s8, 0.3, 5.62, 9.3, 1.6, color=CARD, border=BORDER, rnd=True)
T(s8, 'Why does EfficientNetB2 beat the 4-model ensemble by +9.57%?',
  0.5, 5.68, 9.0, 0.38, sz=12.5, bold=True, color=WHITE)
T(s8, 'ImageNet pretraining exposes EfficientNetB2 to 14M images. The SE-CNN ensemble '
     'trains from scratch on ~84K. Pretrained texture/edge/shape detectors transfer directly '
     'to facial analysis. Architecture design alone cannot replicate this at academic dataset scales.',
  0.5, 6.08, 9.0, 1.0, sz=11, color=GRAY, wrap=True)
NOTES(s8, "Final numbers: 80.02% EfficientNetB2+TTA, +27.22% over baseline, +9.57% over ensemble. "
     "The core lesson: 14M ImageNet pretraining images provide feature richness that a 3.2M "
     "parameter network trained on 84K samples cannot match. Transfer learning wins.")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — PER-CLASS ANALYSIS
# ═══════════════════════════════════════════════════════════════════
s9 = S()
HDR(s9, 'Per-Class Performance: Not All Emotions Are Equal', '80% overall — F1 ranges from 0.91 to 0.42')
IMG(s9, PLOTS/'pptx_perclass_f1.png', 0.3, 0.95, 8.5)
R(s9, 9.05, 0.95, 4.0, 6.35, color=CARD, border=BORDER, rnd=True)
T(s9, 'Key Insights', 9.22, 1.05, 3.65, 0.38, sz=14, bold=True, color=BLUE)
R(s9, 9.05, 1.43, 4.0, 0.03, color=BORDER)
ins = [(GREEN,'Happy  F1=0.91','Duchenne smile = large, distinctive, universal. Easiest.'),
       (RED,'Fear  F1=0.55','Shares muscles with Surprise. 50%+ misclassified.'),
       (RED,'Disgust  F1=0.42','Rarest class. Subtle nasal wrinkle. Annotation noise.'),
       (ORANGE,'Angry  P=0.80 / R=0.64','Conservative: anger only if signal is strong.')]
for i,(col,title,body) in enumerate(ins):
    y=1.55+i*1.3
    R(s9, 9.22, y, 3.65, 1.18, color=CARD2, border=col, bw=0.75, rnd=True)
    T(s9, title, 9.36, y+0.07, 3.38, 0.32, sz=11.5, bold=True, color=col)
    T(s9, body,  9.36, y+0.42, 3.38, 0.68, sz=10,   color=WHITE, wrap=True)
R(s9, 0.3, 6.45, 8.5, 0.82, color=CARD2, border=RED, bw=1.0, rnd=True)
T(s9, '80% overall accuracy coexists with F1=0.42 for disgust and F1=0.55 for fear. '
     'Per-class metrics are mandatory before any real-world deployment.',
  0.5, 6.52, 8.2, 0.65, sz=12, bold=True, color=WHITE, wrap=True)
NOTES(s9, "Per-class analysis is the most important slide for understanding real system quality. "
     "Happy: F1=0.91, large unambiguous deformation. Fear: F1=0.55, shares same muscles as "
     "Surprise, 50%+ misclassified. Disgust: F1=0.42, rarest class, subtle nasal wrinkle, "
     "high annotation ambiguity. The 80% headline hides critical class-level failures.")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — DEPLOYMENT
# ═══════════════════════════════════════════════════════════════════
s10 = S()
HDR(s10, 'Live System: Real-Time Deployment', 'Fully operational — training curves, live webcam, Streamlit, edge export')
IMG(s10, PLOTS/'figure5_training_curves.png', 0.3, 0.95, 8.1)
R(s10, 8.65, 0.95, 4.35, 6.35, color=CARD, border=BORDER, rnd=True)
T(s10, 'Deployment Modes', 8.82, 1.05, 4.0, 0.38, sz=13, bold=True, color=BLUE)
R(s10, 8.65, 1.43, 4.35, 0.03, color=BORDER)
modes = [('Streamlit Web App','5 tabs: Live Demo, Video, Model Comparison, Training Curves, About',BLUE),
         ('Real-Time Webcam','Grad-CAM overlay + prob bars. --compare mode vs DeepFace',CYAN),
         ('CLI Tools','detect_image.py  detect_video.py  meeting_analyzer.py -> PDF',GRAY),
         ('TFLite INT8 Export','4x smaller. Android / Raspberry Pi. No code change.',GREEN)]
for i,(label,desc,col) in enumerate(modes):
    y=1.52+i*1.22
    R(s10, 8.82, y, 4.0, 1.1, color=CARD2, border=col, bw=0.75, rnd=True)
    T(s10, label, 8.98, y+0.06, 3.7, 0.36, sz=11.5, bold=True, color=col)
    T(s10, desc,  8.98, y+0.46, 3.7, 0.55, sz=10,   color=WHITE, wrap=True)
fig4 = PLOTS/'figure4_deployment.png'
if fig4.exists():
    IMG(s10, fig4, 0.3, 6.62, 8.1, 0.7)
else:
    R(s10, 0.3, 6.62, 8.1, 0.65, color=CARD2, border=BLUE, bw=0.75, rnd=True)
    T(s10, '[ Run make_figure4_grid.py to embed deployment screenshots here ]',
      0.5, 6.68, 7.8, 0.45, sz=10, color=GRAY, align=PP_ALIGN.CENTER)
NOTES(s10, "Training curves confirm real experimental behaviour: MobileNetV2 improves across two phases; "
     "EfficientNetB2 shows three phase transitions with Phase 2 producing the steepest gain. "
     "Deployment stack: Streamlit 5-tab app, real-time webcam with Grad-CAM, CLI tools, "
     "meeting analyzer PDF, and TFLite INT8 edge export — all running on a laptop GPU.")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — PROBLEMS AND SOLUTIONS
# ═══════════════════════════════════════════════════════════════════
s11 = S()
HDR(s11, 'Critical Engineering Challenges', 'What went wrong — root cause — what we did')
probs = [
    (RED,   'Augmentation Collapse',
     'RandomBrightness + standardization','Accuracy frozen at 14.3% for 100 epochs',
     'Geometric augmentation only','Preprocessing and augmentation interact — test first'),
    (ORANGE,'Catastrophic Forgetting',
     'High LR fine-tuning overwrites pretrained features','Backbone representations destroyed',
     '4-phase progressive unfreezing at decreasing LR','Warm up task layers before backbone'),
    (BLUE,  'Fear vs Surprise Confusion',
     'Both share frontalis + levator palpebrae','50%+ of fear misclassified as surprise',
     'TTA + ensemble reduce but cannot eliminate','Still images insufficient — need temporal/audio'),
]
for i,(col,title,cause,effect,solution,lesson) in enumerate(probs):
    cx=0.35+i*4.32
    R(s11, cx, 1.0, 4.15, 6.25, color=CARD, border=col, bw=1.5, rnd=True)
    R(s11, cx, 1.0, 4.15, 0.44, color=CARD2)
    T(s11, title, cx+0.14, 1.04, 3.85, 0.34, sz=13, bold=True, color=col)
    for j,(lbl,txt,tc) in enumerate([('Cause:',cause,GRAY),('Effect:',effect,WHITE),('Fix:',solution,GREEN),('Lesson:',lesson,col)]):
        y=1.55+j*1.12
        T(s11, lbl,  cx+0.14, y,      1.05, 0.3,  sz=10, bold=True, color=GRAY)
        T(s11, txt,  cx+0.14, y+0.3,  3.85, 0.75, sz=10.5, color=tc, wrap=True)
        if j<3: R(s11, cx+0.14, y+1.06, 3.85, 0.02, color=BORDER)
NOTES(s11, "Three key challenges: 1) Augmentation collapse — photometric augment + standardization "
     "froze training at 14.3%. Fix: geometric only. 2) Catastrophic forgetting — high LR fine-tuning "
     "destroys pretrained features. Fix: progressive unfreezing. 3) Fear/surprise confusion — "
     "same muscles, needs temporal modeling or audio fusion to properly solve.")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — KEY LESSONS
# ═══════════════════════════════════════════════════════════════════
s12 = S()
HDR(s12, '5 Key Lessons Learned', 'Generalizable insights beyond this project')
lessons = [
    (GOLD,  '01','Pretrained features beat custom architectures at academic scale.',
     'EfficientNetB2 + ImageNet beats 4-model SE-CNN ensemble by +9.57%.'),
    (BLUE,  '02','Preprocessing and augmentation interact — always test the full pipeline first.',
     'Photometric augment + standardisation = training collapse. Geometric only is safe.'),
    (GREEN, '03','Progressive fine-tuning is essential — never fine-tune all weights at once.',
     'Head -> Top 40% -> Full backbone -> SWA. The warm-up phase is not optional.'),
    (ORANGE,'04','Overall accuracy is insufficient for imbalanced FER.',
     'F1=0.42 for disgust hides behind 80% headline. Per-class analysis is mandatory.'),
    (CYAN,  '05','SWA and TTA are free performance boosts — always include in evaluation.',
     'SWA: +0.5-1% by averaging flat-minima weights. TTA: +1-5% at zero training cost.'),
]
for i,(col,num,title,body) in enumerate(lessons):
    y=1.05+i*1.22
    R(s12, 0.35, y, 12.63, 1.1, color=CARD, border=col, bw=1.0, rnd=True)
    R(s12, 0.35, y, 1.1,   1.1, color=CARD2)
    T(s12, num,   0.35, y+0.2,  1.1,  0.65, sz=26, bold=True, color=col, align=PP_ALIGN.CENTER)
    T(s12, title, 1.62, y+0.06, 11.0, 0.38, sz=13.5, bold=True, color=WHITE)
    T(s12, body,  1.62, y+0.52, 11.0, 0.48, sz=11,   color=GRAY, wrap=True)
NOTES(s12, "Five generalizable lessons: 1) Strong pretrained backbone first, then custom design. "
     "2) Validate full preprocessing+augmentation pipeline on a short run before committing. "
     "3) Progressive fine-tuning — head warm-up is not optional. "
     "4) Report per-class metrics, not just overall accuracy. "
     "5) SWA and TTA are essentially free — always include them.")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════
s13 = S()
R(s13, 0, 0, 0.1,   7.5,  color=BLUE)
R(s13, 0, 7.3, 13.33, 0.2, color=BLUE)
T(s13, '52.80%', 0.4,  0.5, 4.5, 1.3, sz=54, bold=True, color=LGRAY, align=PP_ALIGN.CENTER)
T(s13, 'to',     4.85, 0.65, 1.3, 0.95, sz=32, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
T(s13, '80.02%', 6.0,  0.5, 6.8, 1.3, sz=54, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
T(s13, 'DeepFace zero-shot baseline', 0.4, 1.75, 4.5, 0.4, sz=12, color=GRAY, align=PP_ALIGN.CENTER)
T(s13, '+27.22%  |  6 configurations evaluated  |  RAF-DB 3,068 test images',
  0.4, 1.75, 12.5, 0.4, sz=12, color=GRAY, align=PP_ALIGN.CENTER)
T(s13, 'EfficientNetB2 best model', 6.0, 1.75, 6.8, 0.4, sz=12, color=GRAY, align=PP_ALIGN.CENTER)
R(s13, 0.4, 2.28, 12.5, 0.04, color=BLUE)
tks = [
    ('1','ImageNet pretraining decisive — +9.57% over 4-model custom ensemble.',          GOLD),
    ('2','Augmentation must be compatible with preprocessing — photometric -> collapse.',  BLUE),
    ('3','Per-class metrics mandatory — disgust F1=0.42 behind 80% headline.',            RED),
    ('4','Progressive fine-tuning prevents catastrophic forgetting — warm-up first.',     CYAN),
    ('5','Full deployment achieved on laptop GPU — Streamlit, webcam, CLI, TFLite.',      GREEN),
]
for i,(num,text,col) in enumerate(tks):
    y=2.45+i*0.8
    R(s13, 0.4, y, 0.65, 0.65, color=CARD2, border=col, bw=0.75, rnd=True)
    T(s13, num,  0.42, y+0.08, 0.62, 0.48, sz=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    T(s13, text, 1.18, y+0.12, 11.6, 0.48, sz=13, color=WHITE)
R(s13, 0.4, 6.6, 12.5, 0.72, color=CARD2, border=BLUE, bw=0.75, rnd=True)
T(s13, 'Thank you — Questions welcome',
  0.6, 6.65, 12.1, 0.38, sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(s13, 'github.com/charbelmezeraani0-star/facial-emotion-recognition',
  0.6, 7.0,  12.1, 0.28, sz=10, color=GRAY,  align=PP_ALIGN.CENTER)
NOTES(s13, "To conclude: 52.80% to 80.02%, +27.22% improvement across 6 model configurations. "
     "Five takeaways: ImageNet pretrain wins. Preprocessing and augmentation interact. "
     "Per-class metrics are mandatory. Progressive fine-tuning prevents forgetting. "
     "Full deployment is achievable on a laptop GPU. Thank you, happy to take questions.")

# ═══════════════════════════════════════════════════════════════════
prs.save(str(OUT))
print(f'Saved {len(prs.slides)} slides -> {OUT}')
print(f'Size: {OUT.stat().st_size/1024:.0f} KB')
